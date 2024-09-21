import fitz  # PyMuPDF for PDFs
import docx  # python-docx for DOCX
import openpyxl  # openpyxl for Excel files
import csv  # CSV parsing
from pptx import Presentation  # python-pptx for PPTX

# Function to parse different file types
def parse_file(file_path):
    if file_path.endswith(".pdf"):
        return parse_pdf(file_path)
    elif file_path.endswith(".docx"):
        return parse_docx(file_path)
    elif file_path.endswith(".pptx"):
        return parse_pptx(file_path)
    elif file_path.endswith(".xlsx"):
        return parse_xlsx(file_path)
    elif file_path.endswith(".csv"):
        return parse_csv(file_path)
    else:
        raise ValueError(f"Unsupported file type: {file_path}")

# Function to parse PDF files
def parse_pdf(file_path):
    doc = fitz.open(file_path)
    text = ""
    for page in doc:
        text += page.get_text()
    return text

# Function to parse DOCX files
def parse_docx(file_path):
    doc = docx.Document(file_path)
    return "\n".join([paragraph.text for paragraph in doc.paragraphs])

# Function to parse PPTX files
def parse_pptx(file_path):
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

# Function to parse XLSX files
def parse_xlsx(file_path):
    wb = openpyxl.load_workbook(file_path)
    sheet = wb.active
    text = ""
    for row in sheet.iter_rows(values_only=True):
        text += "\t".join([str(cell) if cell else "" for cell in row]) + "\n"
    return text

# Function to parse CSV files
def parse_csv(file_path):
    text = ""
    with open(file_path, newline='', encoding='utf-8') as csvfile:
        reader = csv.reader(csvfile)
        for row in reader:
            text += ",".join(row) + "\n"
    return text
